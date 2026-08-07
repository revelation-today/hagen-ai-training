#!/usr/bin/env python3
"""
Build PowerPoint decks from the session slide outlines.

Reads   : session-NN-<slug>/slides/outline.md
Writes  : decks/NN-<slug>.pptx           (one deck per session)
          decks/mermaid/NN-<slug>/*.mmd  (diagram sources, for later rendering)

Implements the design system in powerpoint_instructions.md:
 - 16:9, one sans family, headline >= 28pt, body >= 18pt
 - neutral palette, no information by colour alone
 - speaker notes in the Notes pane, never on the slide
 - source/licence footer tag on every derived slide
 - alt text on every placeholder

Mermaid note: no Mermaid CLI is available in this environment, so diagrams are
NOT rasterised. Each diagram becomes a labelled placeholder on the slide, and
its source is written both to the speaker notes and to decks/mermaid/. Render
them later with `mmdc -i file.mmd -o file.png` (or mermaid.live) and drop the
images onto the placeholders.

Usage:  python build_decks.py
"""

import re
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------- palette ---
PRIMARY = RGBColor(0x1F, 0x2A, 0x37)   # near-black slate  — headlines
ACCENT = RGBColor(0x25, 0x63, 0xEB)    # blue              — rules, emphasis
MUTED = RGBColor(0x6B, 0x72, 0x80)     # grey              — footers, meta
PANEL = RGBColor(0xF3, 0xF4, 0xF6)     # light grey        — placeholder fill
BODY = RGBColor(0x11, 0x18, 0x27)      # body text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"          # widely available; swap for the corporate standard
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)      # 16:9

ROOT = pathlib.Path(__file__).parent
OUT = ROOT / "decks"
MMD = OUT / "mermaid"

BLOCK_OF = {
    1: "Understand it", 2: "Understand it",
    3: "Methods", 4: "Methods", 5: "Methods", 6: "Methods",
    7: "Do it", 8: "Do it", 9: "Do it",
    10: "Use it well", 11: "Use it well", 12: "Use it well",
    13: "Use it safely", 14: "Use it safely",
    15: "Judge it", 16: "Judge it",
}

# ----------------------------------------------------------------- parsing ---
# The 16 sessions were authored independently and use several label variants:
#   "- **On-slide text:**"   "- **On-slide (headline):**"   "- **On-slide headline:**"
#   "- **Bullets:**"         "- **Visual (Mermaid):**"      "**Speaker notes:**" (no dash)
#   "- **Source/licence:**"  "- **Visual (render this Mermaid):**"
# Match a bold label loosely, then normalise it to a canonical key.
FIELD_RE = re.compile(
    r"^\s*[-*]?\s*\*\*\s*(On-slide[^*:]*|Bullets|Speaker notes|Notes|Visual[^*:]*|"
    r"Source[^*:]*|Licence[^*:]*)\s*:?\s*\*\*\s*:?\s*",
    re.I,
)


def canon_key(label: str) -> str:
    l = label.strip().lower()
    if l.startswith("on-slide") or l == "bullets":
        return "on-slide text"
    if l.startswith("speaker") or l == "notes":
        return "speaker notes"
    if l.startswith("visual"):
        return "visual"
    if l.startswith("source") or l.startswith("licence"):
        return "source/licence"
    return l


def parse_outline(path: pathlib.Path):
    """Split an outline.md into a list of slide dicts."""
    text = path.read_text(encoding="utf-8")
    parts = re.split(r"^#{2,4}\s+Slide\s+", text, flags=re.M)[1:]
    slides = []
    for raw in parts:
        lines = raw.split("\n")
        heading = lines[0].strip()
        m = re.match(r"(\d+)\s*[—\-–:.]\s*(.*)", heading)
        num, title = (m.group(1), m.group(2).strip()) if m else ("", heading)

        fields, current = {}, None
        for ln in lines[1:]:
            fm = FIELD_RE.match(ln)
            if fm:
                key = canon_key(fm.group(1))
                rest = FIELD_RE.sub("", ln).strip()
                # multiple bullets/on-slide fields accumulate rather than overwrite
                fields[key] = (fields.get(key, "") + " · " + rest).strip(" ·") \
                    if key == "on-slide text" and fields.get(key) else rest
                current = key
            elif current:
                fields[current] += "\n" + ln
        # Mermaid is extracted from the WHOLE block: some sessions place the
        # fenced diagram outside any field, after a blank line.
        slides.append({"num": num, "title": title, "raw": raw, **fields})
    return slides


def split_bullets(s: str):
    """On-slide text uses ' · ' as the bullet separator; fall back to sentences."""
    s = re.sub(r"```.*?```", "", s, flags=re.S).strip()
    if not s:
        return []
    if "·" in s:
        items = [b.strip(" .") for b in s.split("·")]
    else:
        items = [b.strip() for b in re.split(r"(?<=[.!?])\s+(?=[A-Z“\"*])", s)]
    return [clean_md(b) for b in items if b.strip()][:6]      # spec: <= 6 bullets


def clean_md(s: str) -> str:
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"\*(.+?)\*", r"\1", s)
    s = re.sub(r"`(.+?)`", r"\1", s)
    s = re.sub(r"\[(.+?)\]\(.*?\)", r"\1", s)
    return re.sub(r"\s+", " ", s).strip()


def extract_mermaid(s: str):
    return re.findall(r"```mermaid\s*(.*?)```", s or "", flags=re.S)


# ----------------------------------------------------------------- drawing ---
def textbox(slide, l, t, w, h, text, size, color=BODY, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=10):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT
    return tb


def bullets(slide, l, t, w, h, items, size=20):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        dot = p.add_run()
        dot.text = "▪  "
        dot.font.size = Pt(size)
        dot.font.color.rgb = ACCENT
        dot.font.name = FONT
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = BODY
        r.font.name = FONT
    return tb


def accent_rule(slide, t=Inches(1.28), l=Inches(0.7), w=Inches(1.1)):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(4))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def footer(slide, session_no, source):
    left = (f"Session {session_no}  ·  AI Training Series" if session_no
            else "AI Training Series")
    textbox(slide, Inches(0.7), Inches(6.92), Inches(7.5), Inches(0.4),
            left, 11, MUTED)
    # Only tag slides whose content is actually *derived* from a source. Clean
    # first, then test — "original." and "none." must not reach the footer.
    tag = clean_md(source or "")
    if re.match(r"^(none|original|n/?a)\b", tag, re.I) or not tag:
        tag = ""
    if tag:
        tag = (tag[:105] + "…") if len(tag) > 106 else tag
        textbox(slide, Inches(8.2), Inches(6.92), Inches(4.4), Inches(0.4),
                tag, 10, MUTED, align=PP_ALIGN.RIGHT)
    return


def diagram_placeholder(slide, l, t, w, h, caption, mermaid_src):
    """A labelled box standing in for an unrendered Mermaid diagram."""
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = MUTED
    box.line.width = Pt(1)
    box.shadow.inherit = False
    box.text_frame.text = ""

    nodes = re.findall(r'\["?(.*?)"?\]', mermaid_src)[:6]
    preview = "  →  ".join(n.replace("<br/>", " ") for n in nodes) if nodes else ""

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "DIAGRAM — render from Mermaid source (in speaker notes)"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = MUTED
    r.font.name = FONT
    if preview:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = preview
        r2.font.size = Pt(15)
        r2.font.color.rgb = PRIMARY
        r2.font.name = FONT
    # accessibility: alt text
    try:
        box._element._nvXxPr.cNvPr.set("descr", f"Diagram: {caption or 'see speaker notes'}")
    except Exception:
        pass
    return box


def set_notes(slide, *chunks):
    body = "\n\n".join(c.strip() for c in chunks if c and c.strip())
    slide.notes_slide.notes_text_frame.text = body


# ------------------------------------------------------------------ build ---
def build_deck(session_dir: pathlib.Path):
    m = re.match(r"session-(\d+)-(.*)", session_dir.name)
    if m:
        no, slug = int(m.group(1)), m.group(2)
    else:
        # standalone decks (e.g. overview-30min): no session number
        no, slug = 0, session_dir.name
    outline = session_dir / "slides" / "outline.md"
    if not outline.exists():
        return None
    slides = parse_outline(outline)

    # session title from the session README's H1
    readme = (session_dir / "README.md").read_text(encoding="utf-8")
    st = re.search(r"^#\s*(.+)$", readme, re.M)
    session_title = re.sub(r"^Session\s+\d+\s+[—-]\s+", "", st.group(1)).strip() if st else slug

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    mmd_dir = MMD / (f"{no:02d}-{slug}" if no else slug)
    diagrams = 0

    for i, s in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        title = clean_md(s["title"])
        on_slide = s.get("on-slide text", "")
        notes = clean_md_notes(s.get("speaker notes", ""))
        visual = s.get("visual", "")
        source = s.get("source/licence", "")
        mermaids = extract_mermaid(s.get("raw", "")) or extract_mermaid(visual)

        is_title = i == 0 and re.match(r"title", title, re.I)

        if is_title:
            textbox(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.6),
                    session_title, 40, PRIMARY, bold=True)
            accent_rule(slide, t=Inches(4.05), l=Inches(0.9))
            subtitle = (f"Session {no} of 16  ·  {BLOCK_OF.get(no,'')}"
                        if no else "Standalone intro talk")
            duration = ("AI Training Series  ·  45 min + 15 min Q&A" if no
                        else "AI Training Series  ·  ~30 min + Q&A")
            textbox(slide, Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.9),
                    [subtitle, duration], 17, MUTED)
        else:
            textbox(slide, Inches(0.7), Inches(0.55), Inches(11.9), Inches(0.8),
                    title, 28, PRIMARY, bold=True)
            accent_rule(slide)

            items = split_bullets(on_slide)
            if mermaids:
                if items:
                    bullets(slide, Inches(0.7), Inches(1.6), Inches(5.4),
                            Inches(4.6), items, size=17)
                    diagram_placeholder(slide, Inches(6.5), Inches(1.7),
                                        Inches(6.1), Inches(4.2), title, mermaids[0])
                else:
                    diagram_placeholder(slide, Inches(1.4), Inches(1.9),
                                        Inches(10.5), Inches(4.2), title, mermaids[0])
            elif items:
                bullets(slide, Inches(0.9), Inches(1.85), Inches(11.4),
                        Inches(4.4), items, size=21)
            else:
                textbox(slide, Inches(0.9), Inches(2.4), Inches(11.4), Inches(2.0),
                        clean_md(on_slide)[:400], 20, BODY)

            footer(slide, no, source)

        # notes: narration + any mermaid source for the designer
        extra = ""
        for j, mm in enumerate(mermaids):
            extra += f"\n\n--- Mermaid source {j+1} (render and place on the diagram box) ---\n{mm.strip()}"
            mmd_dir.mkdir(parents=True, exist_ok=True)
            (mmd_dir / f"slide-{s['num'] or i:0>2}-{j+1}.mmd").write_text(mm.strip(), encoding="utf-8")
            diagrams += 1
        set_notes(slide, notes, extra)

    OUT.mkdir(parents=True, exist_ok=True)
    path = OUT / (f"{no:02d}-{slug}.pptx" if no else f"{slug}.pptx")
    prs.save(path)
    return path, len(slides), diagrams


def clean_md_notes(s: str) -> str:
    s = re.sub(r"```.*?```", "", s or "", flags=re.S)
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"\*(.+?)\*", r"\1", s)
    s = re.sub(r"`(.+?)`", r"\1", s)
    return re.sub(r"\n{3,}", "\n\n", s).strip()


if __name__ == "__main__":
    dirs = sorted(ROOT.glob("session-*")) + sorted(ROOT.glob("overview-*"))
    total_slides = total_diagrams = 0
    print(f"{'deck':<34} {'slides':>7} {'diagrams':>9}")
    print("-" * 53)
    for d in dirs:
        res = build_deck(d)
        if res:
            path, n, dg = res
            total_slides += n
            total_diagrams += dg
            print(f"{path.name:<34} {n:>7} {dg:>9}")
    print("-" * 53)
    print(f"{'TOTAL':<34} {total_slides:>7} {total_diagrams:>9}")
    print(f"\nDecks   : {OUT}")
    print(f"Diagrams: {MMD}  (render with: mmdc -i x.mmd -o x.png)")
